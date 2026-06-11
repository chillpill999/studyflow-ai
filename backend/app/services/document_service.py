import os
import re
import time
from typing import List, Dict, Any
import PyPDF2
import docx2txt
from pptx import Presentation

class DocumentService:
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> str:
        """
        Extract raw text from a document based on its extension/type.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_type = file_type.lower()
        text = ""

        if file_type == "pdf":
            try:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                
                # Fallback to Gemini OCR for scanned PDFs
                if len(text.strip()) < 50:
                    from app.core.config import settings
                    if settings.GEMINI_API_KEY:
                        import google.generativeai as genai
                        genai.configure(api_key=settings.GEMINI_API_KEY)
                        uploaded_file = genai.upload_file(file_path)
                        
                        # Wait for processing
                        while uploaded_file.state.name == "PROCESSING":
                            time.sleep(2)
                            uploaded_file = genai.get_file(uploaded_file.name)
                            
                        if uploaded_file.state.name == "FAILED":
                            raise ValueError("Gemini File processing failed.")
                            
                        model = genai.GenerativeModel("gemini-1.5-flash")
                        response = model.generate_content([
                            uploaded_file, 
                            "Extract all the text from this document accurately. Respond ONLY with the raw extracted text."
                        ])
                        text = response.text
                        genai.delete_file(uploaded_file.name)
            except Exception as e:
                raise ValueError(f"Failed to parse PDF: {str(e)}")

        elif file_type == "docx":
            try:
                text = docx2txt.process(file_path)
            except Exception as e:
                raise ValueError(f"Failed to parse DOCX: {str(e)}")

        elif file_type == "pptx":
            try:
                prs = Presentation(file_path)
                slide_texts = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text.append(shape.text)
                    slide_texts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
                text = "\n\n".join(slide_texts)
            except Exception as e:
                raise ValueError(f"Failed to parse PPTX: {str(e)}")

        elif file_type in ["txt", "text", "md"]:
            try:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            except Exception as e:
                raise ValueError(f"Failed to parse text file: {str(e)}")
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        # Basic text cleaning: remove extra whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 800, chunk_overlap: int = 150) -> List[Dict[str, Any]]:
        """
        Splits text into overlapping chunks. Returns list of dicts with 'text' and 'index'.
        """
        if not text:
            return []

        chunks = []
        words = text.split(" ")
        
        current_chunk_words = []
        current_length = 0
        chunk_idx = 0

        for word in words:
            current_chunk_words.append(word)
            current_length += len(word) + 1  # count space
            
            if current_length >= chunk_size:
                chunk_content = " ".join(current_chunk_words)
                chunks.append({
                    "id": chunk_idx,
                    "text": chunk_content
                })
                chunk_idx += 1
                
                # Keep overlap: approximate words for overlap
                # Keep the last 15-20 words as overlap
                overlap_words = current_chunk_words[-20:] if len(current_chunk_words) > 20 else current_chunk_words[-5:]
                current_chunk_words = list(overlap_words)
                current_length = sum(len(w) + 1 for w in current_chunk_words)

        if current_chunk_words:
            chunk_content = " ".join(current_chunk_words)
            if chunk_content.strip():
                chunks.append({
                    "id": chunk_idx,
                    "text": chunk_content
                })

        return chunks
